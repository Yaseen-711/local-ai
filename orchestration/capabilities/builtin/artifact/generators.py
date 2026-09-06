"""Deterministic document format generators (XLSX, DOCX, PDF)."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Union

from orchestration.capabilities.builtin.artifact.types import ArtifactGenerationRequest


def _normalize_table_data(
    raw_data: Union[List[List[Any]], List[Dict[str, Any]], Dict[str, Any], None]
) -> Dict[str, List[List[str]]]:
    """Normalize raw structured data into a mapping of {sheet_name: 2D string grid}."""
    if raw_data is None:
        return {"Sheet1": []}

    # Case 1: Dict of sheets { "SheetName": [[...], [...]] }
    if isinstance(raw_data, dict):
        # Check if it looks like a single dict record or a multi-sheet mapping
        first_val = next(iter(raw_data.values()), None) if raw_data else None
        if isinstance(first_val, list):
            result = {}
            for sheet_name, sheet_data in raw_data.items():
                norm = _normalize_table_data(sheet_data)
                result[sheet_name] = norm.get("Sheet1", [])
            return result
        else:
            # Single dict record -> keys as headers, values as row
            headers = list(raw_data.keys())
            row = [str(raw_data[k]) for k in headers]
            return {"Sheet1": [headers, row]}

    # Case 2: List of dicts (standard tabular format)
    if isinstance(raw_data, list) and len(raw_data) > 0 and isinstance(raw_data[0], dict):
        headers = list(raw_data[0].keys())
        rows = [[item.get(h, "") for h in headers] for item in raw_data]
        return {"Sheet1": [headers] + rows}

    # Case 3: 2D List of rows
    if isinstance(raw_data, list):
        rows = [[cell for cell in row] if isinstance(row, list) else [row] for row in raw_data]
        return {"Sheet1": rows}

    return {"Sheet1": [[raw_data]]}



class XlsxGenerator:
    """Deterministic XLSX spreadsheet generator using openpyxl."""

    @staticmethod
    def generate(request: ArtifactGenerationRequest, output_path: Path) -> None:
        import openpyxl
        from openpyxl.styles import Alignment, Font, PatternFill

        wb = openpyxl.Workbook()
        # Remove default active sheet if we create custom sheets
        wb.remove(wb.active)  # type: ignore[arg-type]

        sheets = _normalize_table_data(request.data)
        if request.title and list(sheets.keys()) == ["Sheet1"]:
            sanitized_title = request.title.replace(":", " ").replace("/", " ").replace("\\", " ")[:31].strip() or "Sheet1"
            sheets = {sanitized_title: sheets["Sheet1"]}

        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        cell_font = Font(name="Calibri", size=10)

        for sheet_name, grid in sheets.items():
            ws = wb.create_sheet(title=sheet_name[:31])
            ws.views.sheetView[0].showGridLines = True

            for r_idx, row in enumerate(grid, start=1):
                for c_idx, val in enumerate(row, start=1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=val)
                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal="center", vertical="center")
                    else:
                        cell.font = cell_font
                        cell.alignment = Alignment(vertical="center")

            # Auto-fit column widths
            for col in ws.columns:
                max_len = max(len(str(cell.value or "")) for cell in col)
                col_letter = col[0].column_letter
                ws.column_dimensions[col_letter].width = max(max_len + 3, 10)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(output_path))


class DocxGenerator:
    """Deterministic DOCX document generator using python-docx."""

    @staticmethod
    def generate(request: ArtifactGenerationRequest, output_path: Path) -> None:
        import docx

        doc = docx.Document()

        if request.title:
            doc.add_heading(request.title, level=0)

        if request.content:
            lines = [line.rstrip() for line in request.content.splitlines()]
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if not line:
                    i += 1
                    continue
                if line.startswith("# "):
                    doc.add_heading(line[2:].strip(), level=1)
                    i += 1
                elif line.startswith("## "):
                    doc.add_heading(line[3:].strip(), level=2)
                    i += 1
                elif line.startswith("### "):
                    doc.add_heading(line[4:].strip(), level=3)
                    i += 1
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line.lstrip("-* ").strip(), style="List Bullet")
                    i += 1
                else:
                    para_lines = [line]
                    i += 1
                    while i < len(lines) and lines[i].strip() and not lines[i].strip().startswith(("#", "-", "*")):
                        para_lines.append(lines[i].strip())
                        i += 1
                    doc.add_paragraph(" ".join(para_lines))

        sheets = _normalize_table_data(request.data)
        for sheet_name, grid in sheets.items():
            if grid:
                if len(sheets) > 1:
                    doc.add_heading(sheet_name, level=2)

                num_rows = len(grid)
                num_cols = len(grid[0]) if num_rows > 0 else 0
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = "Table Grid"

                for r_idx, row in enumerate(grid):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = str(val)
                        if r_idx == 0:
                            for p in cell.paragraphs:
                                for run in p.runs:
                                    run.font.bold = True

        doc.save(str(output_path))


class PdfGenerator:
    """Deterministic PDF document generator using reportlab."""

    @staticmethod
    def generate(request: ArtifactGenerationRequest, output_path: Path) -> None:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=letter,
            leftMargin=54,
            rightMargin=54,
            topMargin=54,
            bottomMargin=54,
        )

        styles = getSampleStyleSheet()
        story = []

        if request.title:
            story.append(Paragraph(request.title, styles["Title"]))
            story.append(Spacer(1, 14))

        if request.content:
            for block in request.content.split("\n\n"):
                b_clean = block.strip()
                if b_clean:
                    if b_clean.startswith("# "):
                        story.append(Paragraph(b_clean[2:], styles["Heading1"]))
                    elif b_clean.startswith("## "):
                        story.append(Paragraph(b_clean[3:], styles["Heading2"]))
                    else:
                        story.append(Paragraph(b_clean, styles["Normal"]))
                    story.append(Spacer(1, 8))

        sheets = _normalize_table_data(request.data)
        for sheet_name, grid in sheets.items():
            if grid:
                if len(sheets) > 1:
                    story.append(Paragraph(sheet_name, styles["Heading2"]))
                    story.append(Spacer(1, 6))

                # Wrap cell contents in Paragraphs for text wrapping
                cell_style = styles["Normal"]
                t_data = [
                    [Paragraph(str(cell), cell_style) for cell in row]
                    for row in grid
                ]

                pdf_table = Table(t_data)
                pdf_table.setStyle(
                    TableStyle([
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                        ("TOPPADDING", (0, 0), (-1, -1), 6),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CCCCCC")),
                    ])
                )
                story.append(pdf_table)
                story.append(Spacer(1, 14))

        doc.build(story)


class PptxGenerator:
    """Deterministic PPTX presentation generator using python-pptx."""

    @staticmethod
    def generate(request: ArtifactGenerationRequest, output_path: Path) -> None:
        import pptx
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 1. Title Slide if title provided
        if request.title:
            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            title = slide.shapes.title
            if title:
                title.text = request.title

        # 2. Content Slides from markdown prose
        if request.content:
            lines = [line.rstrip() for line in request.content.splitlines()]
            current_heading: str | None = None
            current_bullets: List[str] = []
            current_paragraphs: List[str] = []

            def _flush_content_slide() -> None:
                nonlocal current_heading, current_bullets, current_paragraphs
                if not current_heading and not current_bullets and not current_paragraphs:
                    return
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
                if slide.shapes.title:
                    slide.shapes.title.text = current_heading or "Overview"
                    for p in slide.shapes.title.text_frame.paragraphs:
                        p.font.size = Pt(24)
                        p.font.bold = True

                tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
                tf = tx_box.text_frame
                tf.word_wrap = True

                is_first = True
                for p_text in current_paragraphs:
                    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                    is_first = False
                    p.text = p_text
                    p.font.size = Pt(16)
                    p.space_after = Pt(10)

                for b_text in current_bullets:
                    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                    is_first = False
                    p.text = f"•  {b_text}"
                    p.font.size = Pt(15)
                    p.space_after = Pt(6)

                current_heading = None
                current_bullets = []
                current_paragraphs = []

            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if not line:
                    i += 1
                    continue
                if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
                    _flush_content_slide()
                    current_heading = line.lstrip("#").strip()
                    i += 1
                elif line.startswith("- ") or line.startswith("* "):
                    current_bullets.append(line.lstrip("-* ").strip())
                    i += 1
                else:
                    para_lines = [line]
                    i += 1
                    while i < len(lines) and lines[i].strip() and not lines[i].strip().startswith(("#", "-", "*")):
                        para_lines.append(lines[i].strip())
                        i += 1
                    current_paragraphs.append(" ".join(para_lines))

            _flush_content_slide()

        # 3. Data Tables as Presentation Slides
        sheets = _normalize_table_data(request.data)
        for sheet_name, grid in sheets.items():
            if grid and len(grid) > 0:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                if slide.shapes.title:
                    slide.shapes.title.text = sheet_name
                    for p in slide.shapes.title.text_frame.paragraphs:
                        p.font.size = Pt(24)
                        p.font.bold = True

                num_rows = len(grid)
                num_cols = len(grid[0]) if num_rows > 0 else 0
                if num_cols > 0 and num_rows > 0:
                    tbl_width = Inches(11.7)
                    tbl_height = min(Inches(5.0), Inches(0.4 * num_rows + 0.6))
                    table_shape = slide.shapes.add_table(
                        num_rows, num_cols, Inches(0.8), Inches(1.8), tbl_width, tbl_height
                    )
                    table = table_shape.table

                    for r_idx, row in enumerate(grid):
                        for c_idx, val in enumerate(row):
                            cell = table.cell(r_idx, c_idx)
                            cell.text = str(val)
                            if r_idx == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
                                for p in cell.text_frame.paragraphs:
                                    p.font.bold = True
                                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                    p.font.size = Pt(12)
                            else:
                                for p in cell.text_frame.paragraphs:
                                    p.font.size = Pt(11)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

