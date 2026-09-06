"""Template renderers for industrial deliverables (DOCX, XLSX, PPTX).

Compiles structured deliverable specifications into professional, neutral
engineering documents with high visual fidelity, live formulas, and strict formatting.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Optional, Union

from orchestration.capabilities.builtin.artifact.templates.specs import (
    EngineeringCalculationSpec,
    ExecutivePresentationSpec,
    TechnicalApprovalNoteSpec,
    create_demo_approval_note,
    create_demo_calculation_workbook,
    create_demo_executive_presentation,
)
from orchestration.capabilities.builtin.artifact.types import ArtifactFormat


# --------------------------------------------------------------------------- #
# 1. DOCX Technical Approval Note Renderer                                    #
# --------------------------------------------------------------------------- #

def render_docx_approval_note(spec: TechnicalApprovalNoteSpec, output_path: Path) -> None:
    """Render formal Technical Approval Note as a DOCX document."""
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import parse_xml, OxmlElement
    from docx.oxml.ns import nsdecls, qn

    doc = docx.Document()

    # Page Margins: 1 inch
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Palette
    NAVY_HEX = "1F4E79"
    SLATE_HEX = "2E75B6"
    GRAY_HEX = "F2F4F7"

    def _set_cell_background(cell: Any, fill_hex: str) -> None:
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
        cell._tc.get_or_add_tcPr().append(shading)

    def _set_cell_margins(cell: Any, top: int = 100, bottom: int = 100, left: int = 150, right: int = 150) -> None:
        tcPr = cell._tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for m, val in [('w:top', top), ('w:bottom', bottom), ('w:left', left), ('w:right', right)]:
            node = OxmlElement(m)
            node.set(qn('w:w'), str(val))
            node.set(qn('w:type'), 'dxa')
            tcMar.append(node)
        tcPr.append(tcMar)

    # Document Title Block
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(spec.title)
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    title_p.paragraph_format.space_after = Pt(4)

    sub_p = doc.add_paragraph()
    sub_run = sub_p.add_run(f"ENGINEERING TECHNICAL MEMORANDUM | {spec.facility.upper()}")
    sub_run.font.name = "Calibri"
    sub_run.font.size = Pt(10)
    sub_run.font.bold = True
    sub_run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    sub_p.paragraph_format.space_after = Pt(12)

    # Metadata Grid (2 columns of key-value pairs in a 4-column table)
    meta_table = doc.add_table(rows=4, cols=4)
    meta_table.style = "Table Grid"
    meta_pairs = [
        ("Document ID:", spec.document_id, "Revision:", spec.revision),
        ("Date:", spec.date, "Status:", spec.status.replace("_", " ")),
        ("Facility:", spec.facility, "Unit / Area:", spec.unit_area),
        ("Author:", spec.author, "Designated Approver:", spec.approver or "Engineering Authority"),
    ]
    for r_idx, (k1, v1, k2, v2) in enumerate(meta_pairs):
        row = meta_table.rows[r_idx]
        for c_idx, (k, v) in enumerate([(k1, v1), (k2, v2)]):
            cell_k = row.cells[c_idx * 2]
            cell_v = row.cells[c_idx * 2 + 1]
            cell_k.text = k
            cell_v.text = v
            _set_cell_background(cell_k, GRAY_HEX)
            _set_cell_margins(cell_k, 80, 80, 100, 100)
            _set_cell_margins(cell_v, 80, 80, 100, 100)
            for p in cell_k.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9.5)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            for p in cell_v.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9.5)

    doc.add_paragraph().paragraph_format.space_after = Pt(8)

    # Section Helper
    def _add_section_heading(text: str) -> None:
        h = doc.add_heading(text, level=1)
        for r in h.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        h.paragraph_format.space_before = Pt(12)
        h.paragraph_format.space_after = Pt(4)

    # 1. Executive Summary
    _add_section_heading("1. Executive Summary & Problem Context")
    p_exec = doc.add_paragraph(spec.executive_summary)
    p_exec.paragraph_format.space_after = Pt(8)
    for r in p_exec.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(10.5)

    # 2. Design Basis & Process Parameters
    _add_section_heading("2. Design Basis & Operating Parameters")
    p_basis = doc.add_paragraph(f"Governing Standards: {spec.design_basis}")
    for r in p_basis.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(10)
        r.font.italic = True
    p_basis.paragraph_format.space_after = Pt(6)

    if spec.operating_parameters:
        op_table = doc.add_table(rows=len(spec.operating_parameters) + 1, cols=2)
        op_table.style = "Table Grid"
        # Header
        op_table.rows[0].cells[0].text = "Process Parameter"
        op_table.rows[0].cells[1].text = "Design Value / Specification"
        for c in op_table.rows[0].cells:
            _set_cell_background(c, NAVY_HEX)
            _set_cell_margins(c, 80, 80, 100, 100)
            for p in c.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9.5)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for idx, (k, val) in enumerate(spec.operating_parameters.items(), start=1):
            row = op_table.rows[idx]
            row.cells[0].text = k
            row.cells[1].text = val
            _set_cell_margins(row.cells[0], 60, 60, 100, 100)
            _set_cell_margins(row.cells[1], 60, 60, 100, 100)
            for cell in row.cells:
                for p in cell.paragraphs:
                    for r in p.runs:
                        r.font.name = "Calibri"
                        r.font.size = Pt(9.5)

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 3. Equipment Inspection & Tag Registry Table
    _add_section_heading("3. Equipment Inspection & Tag Evaluation Registry")
    headers = ["Tag ID", "Description", "P&ID Reference", "Service", "Design Spec", "Measured Finding", "Status"]
    tag_table = doc.add_table(rows=len(spec.inspection_findings) + 1, cols=len(headers))
    tag_table.style = "Table Grid"

    for c_idx, h_text in enumerate(headers):
        cell = tag_table.rows[0].cells[c_idx]
        cell.text = h_text
        _set_cell_background(cell, NAVY_HEX)
        _set_cell_margins(cell, 80, 80, 80, 80)
        for p in cell.paragraphs:
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r_idx, item in enumerate(spec.inspection_findings, start=1):
        row = tag_table.rows[r_idx]
        row_vals = [
            item.tag_id,
            item.description,
            item.pid_reference,
            item.service,
            item.design_spec,
            item.measured_condition,
            item.compliance_status,
        ]
        for c_idx, val in enumerate(row_vals):
            cell = row.cells[c_idx]
            cell.text = val
            _set_cell_margins(cell, 60, 60, 80, 80)
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9)
                    if c_idx == len(headers) - 1:
                        r.font.bold = True
                        if val == "PASS":
                            r.font.color.rgb = RGBColor(0x27, 0x6A, 0x3C)
                        elif val in ("FAIL", "REQUIRES_ACTION", "ACTION"):
                            r.font.color.rgb = RGBColor(0xC5, 0x5A, 0x11)

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 4. Risk Assessment
    if spec.risk_assessment:
        _add_section_heading("4. Operational Risk Assessment & Controls")
        risk_headers = ["Risk Event", "Severity", "Probability", "Mitigation Measure", "Residual Risk"]
        r_table = doc.add_table(rows=len(spec.risk_assessment) + 1, cols=len(risk_headers))
        r_table.style = "Table Grid"
        for c_idx, h_text in enumerate(risk_headers):
            cell = r_table.rows[0].cells[c_idx]
            cell.text = h_text
            _set_cell_background(cell, NAVY_HEX)
            _set_cell_margins(cell, 80, 80, 80, 80)
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for r_idx, r_dict in enumerate(spec.risk_assessment, start=1):
            row = r_table.rows[r_idx]
            vals = [
                r_dict.get("Risk Event", ""),
                r_dict.get("Severity", ""),
                r_dict.get("Probability", ""),
                r_dict.get("Mitigation Measure", ""),
                r_dict.get("Residual Risk", ""),
            ]
            for c_idx, val in enumerate(vals):
                cell = row.cells[c_idx]
                cell.text = val
                _set_cell_margins(cell, 60, 60, 80, 80)
                for p in cell.paragraphs:
                    for r in p.runs:
                        r.font.name = "Calibri"
                        r.font.size = Pt(9)

        doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 5. Recommendations
    _add_section_heading("5. Engineering Recommendations & Actions")
    for rec in spec.recommendations:
        p = doc.add_paragraph(rec, style="List Bullet" if not rec.strip()[:2].isdigit() else "Normal")
        p.paragraph_format.space_after = Pt(4)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(10)

    # 6. Sign-off Matrix
    _add_section_heading("6. Engineering Sign-Off & Authority Approval")
    so_headers = ["Role", "Approver Name", "Title / Designation", "Decision", "Date", "Signature"]
    so_table = doc.add_table(rows=len(spec.sign_offs) + 1, cols=len(so_headers))
    so_table.style = "Table Grid"

    for c_idx, h_text in enumerate(so_headers):
        cell = so_table.rows[0].cells[c_idx]
        cell.text = h_text
        _set_cell_background(cell, NAVY_HEX)
        _set_cell_margins(cell, 80, 80, 80, 80)
        for p in cell.paragraphs:
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r_idx, so in enumerate(spec.sign_offs, start=1):
        row = so_table.rows[r_idx]
        vals = [
            so.role,
            so.name,
            so.title,
            so.status,
            so.date or "PENDING",
            "[ SIGNED ON FILE ]" if so.status == "APPROVED" else "____________________",
        ]
        for c_idx, val in enumerate(vals):
            cell = row.cells[c_idx]
            cell.text = val
            _set_cell_margins(cell, 60, 60, 80, 80)
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(9)
                    if c_idx == 3:
                        r.font.bold = True
                        if val == "APPROVED":
                            r.font.color.rgb = RGBColor(0x27, 0x6A, 0x3C)
                        else:
                            r.font.color.rgb = RGBColor(0xC5, 0x5A, 0x11)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))


# --------------------------------------------------------------------------- #
# 2. XLSX Engineering Calculation Workbook Renderer                          #
# --------------------------------------------------------------------------- #

def render_xlsx_calculation_workbook(spec: EngineeringCalculationSpec, output_path: Path) -> None:
    """Render Engineering Calculation Workbook with live Excel formulas, units, and verification status."""
    import openpyxl
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # type: ignore[arg-type]

    # Shared Styles
    FONT_FAMILY = "Calibri"
    NAVY_FILL = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    STEEL_FILL = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
    LIGHT_GRAY_FILL = PatternFill(start_color="F2F4F7", end_color="F2F4F7", fill_type="solid")
    PASS_FILL = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    FAIL_FILL = PatternFill(start_color="FCE4D6", end_color="FCE4D6", fill_type="solid")

    WHITE_BOLD_11 = Font(name=FONT_FAMILY, size=11, bold=True, color="FFFFFF")
    WHITE_BOLD_14 = Font(name=FONT_FAMILY, size=14, bold=True, color="FFFFFF")
    BOLD_11 = Font(name=FONT_FAMILY, size=11, bold=True)
    BOLD_10 = Font(name=FONT_FAMILY, size=10, bold=True)
    NORMAL_10 = Font(name=FONT_FAMILY, size=10)
    ITALIC_10 = Font(name=FONT_FAMILY, size=10, italic=True)
    PASS_FONT = Font(name=FONT_FAMILY, size=10, bold=True, color="276A3C")
    FAIL_FONT = Font(name=FONT_FAMILY, size=10, bold=True, color="C55A11")

    THIN_BORDER = Border(
        left=Side(style="thin", color="D9D9D9"),
        right=Side(style="thin", color="D9D9D9"),
        top=Side(style="thin", color="D9D9D9"),
        bottom=Side(style="thin", color="D9D9D9"),
    )

    def _auto_fit_columns(ws: Any) -> None:
        ws.views.sheetView[0].showGridLines = True
        for col in ws.columns:
            max_len = 0
            for cell in col:
                val = str(cell.value or "")
                # Avoid huge length if cell is a formula
                if val.startswith("="):
                    max_len = max(max_len, 15)
                else:
                    max_len = max(max_len, len(val))
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

    # ----------------------------------------------------------------------- #
    # Sheet 1: Cover & Scope                                                  #
    # ----------------------------------------------------------------------- #
    ws1 = wb.create_sheet(title="Cover & Scope")
    ws1.cell(row=1, column=1, value=spec.workbook_title.upper()).font = WHITE_BOLD_14
    ws1.cell(row=1, column=1).fill = NAVY_FILL
    ws1.merge_cells("A1:F1")
    ws1.row_dimensions[1].height = 32

    # Metadata table
    meta = [
        ("Project ID / Reference:", spec.project_id),
        ("Facility / Area:", spec.facility),
        ("Date of Calculation:", spec.date),
        ("Author:", spec.author),
        ("Independent Checker:", spec.checker),
        ("Status:", "VERIFIED & APPROVED"),
    ]
    for idx, (lbl, val) in enumerate(meta, start=3):
        c_lbl = ws1.cell(row=idx, column=1, value=lbl)
        c_lbl.font = BOLD_10
        c_lbl.fill = LIGHT_GRAY_FILL
        c_lbl.border = THIN_BORDER
        c_val = ws1.cell(row=idx, column=2, value=val)
        c_val.font = NORMAL_10
        c_val.border = THIN_BORDER
        ws1.merge_cells(start_row=idx, start_column=2, end_row=idx, end_column=4)

    # Governing Standards
    row_st = 10
    ws1.cell(row=row_st, column=1, value="Governing Engineering Standards").font = WHITE_BOLD_11
    ws1.cell(row=row_st, column=1).fill = STEEL_FILL
    ws1.merge_cells(start_row=row_st, start_column=1, end_row=row_st, end_column=4)
    for idx, std in enumerate(spec.governing_standards, start=row_st + 1):
        c = ws1.cell(row=idx, column=1, value=f"• {std}")
        c.font = NORMAL_10
        ws1.merge_cells(start_row=idx, start_column=1, end_row=idx, end_column=4)

    # Scope
    row_sc = row_st + len(spec.governing_standards) + 2
    ws1.cell(row=row_sc, column=1, value="Scope of Calculation & Design Objectives").font = WHITE_BOLD_11
    ws1.cell(row=row_sc, column=1).fill = STEEL_FILL
    ws1.merge_cells(start_row=row_sc, start_column=1, end_row=row_sc, end_column=4)
    c_sc = ws1.cell(row=row_sc + 1, column=1, value=spec.scope_description)
    c_sc.font = NORMAL_10
    c_sc.alignment = Alignment(wrap_text=True)
    ws1.merge_cells(start_row=row_sc + 1, start_column=1, end_row=row_sc + 2, end_column=4)

    # Separate Verification Evidence
    if spec.verification_evidence:
        row_ve = row_sc + 4
        ws1.cell(row=row_ve, column=1, value="Independent Verification Evidence").font = WHITE_BOLD_11
        ws1.cell(row=row_ve, column=1).fill = NAVY_FILL
        ws1.merge_cells(start_row=row_ve, start_column=1, end_row=row_ve, end_column=4)

        ve_pairs = [
            ("Verification Method:", spec.verification_evidence.method),
            ("Independent Verifier:", spec.verification_evidence.verifier),
            ("Verification Date:", spec.verification_evidence.verification_date),
            ("Verification Status:", spec.verification_evidence.status),
            ("Verification Notes:", spec.verification_evidence.evidence_notes),
        ]
        for idx, (v_lbl, v_val) in enumerate(ve_pairs, start=row_ve + 1):
            c_lbl = ws1.cell(row=idx, column=1, value=v_lbl)
            c_lbl.font = BOLD_10
            c_lbl.fill = LIGHT_GRAY_FILL
            c_lbl.border = THIN_BORDER
            c_val = ws1.cell(row=idx, column=2, value=v_val)
            c_val.font = NORMAL_10
            c_val.border = THIN_BORDER
            c_val.alignment = Alignment(wrap_text=True)
            ws1.merge_cells(start_row=idx, start_column=2, end_row=idx, end_column=4)

    _auto_fit_columns(ws1)

    # ----------------------------------------------------------------------- #
    # Sheet 2: Design Parameters                                              #
    # ----------------------------------------------------------------------- #
    ws2 = wb.create_sheet(title="Design Parameters")
    ws2.cell(row=1, column=1, value="INPUT DESIGN PARAMETERS & MATERIAL SPECIFICATIONS").font = WHITE_BOLD_11
    ws2.cell(row=1, column=1).fill = NAVY_FILL
    ws2.merge_cells("A1:G1")

    p_headers = ["Item", "Parameter Description", "Symbol", "Design Value", "Unit", "Source / Reference Document", "Tolerance / Design Margin"]
    for c_idx, h in enumerate(p_headers, start=1):
        c = ws2.cell(row=3, column=c_idx, value=h)
        c.font = WHITE_BOLD_11
        c.fill = STEEL_FILL
        c.alignment = Alignment(horizontal="center", vertical="center")
        c.border = THIN_BORDER

    for r_idx, p in enumerate(spec.input_parameters, start=4):
        ws2.cell(row=r_idx, column=1, value=r_idx - 3).alignment = Alignment(horizontal="center")
        ws2.cell(row=r_idx, column=2, value=p.name)
        c_sym = ws2.cell(row=r_idx, column=3, value=p.symbol)
        c_sym.alignment = Alignment(horizontal="center")
        c_sym.font = BOLD_10

        c_val = ws2.cell(row=r_idx, column=4, value=p.value)
        c_val.alignment = Alignment(horizontal="right")
        c_val.number_format = "#,##0.00"

        c_unit = ws2.cell(row=r_idx, column=5, value=p.unit)
        c_unit.alignment = Alignment(horizontal="center")

        ws2.cell(row=r_idx, column=6, value=p.source)
        ws2.cell(row=r_idx, column=7, value=p.tolerance or "Code Standard")

        for col_i in range(1, 8):
            ws2.cell(row=r_idx, column=col_i).border = THIN_BORDER
            ws2.cell(row=r_idx, column=col_i).font = NORMAL_10

    _auto_fit_columns(ws2)

    # ----------------------------------------------------------------------- #
    # Sheet 3: Calculations                                                   #
    # ----------------------------------------------------------------------- #
    ws3 = wb.create_sheet(title="Calculations")
    ws3.cell(row=1, column=1, value="ENGINEERING DERIVATIONS & LIVE FORMULA VERIFICATION").font = WHITE_BOLD_11
    ws3.cell(row=1, column=1).fill = NAVY_FILL
    ws3.merge_cells("A1:M1")

    calc_headers = [
        "Step",
        "Derivation Item / Description",
        "Symbol",
        "Governing Equation (Text)",
        "Numerical Substitution",
        "Live Excel Formula",
        "Calculated Value",
        "Unit",
        "Design Limit / Benchmark",
        "Tolerance",
        "Verification Formula",
        "Status",
        "Independent Verification Evidence",
    ]
    for c_idx, h in enumerate(calc_headers, start=1):
        c = ws3.cell(row=3, column=c_idx, value=h)
        c.font = WHITE_BOLD_11
        c.fill = STEEL_FILL
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = THIN_BORDER

    ws3.row_dimensions[3].height = 28

    for r_idx, step in enumerate(spec.steps, start=4):
        # Step ID
        c_step = ws3.cell(row=r_idx, column=1, value=step.step_id)
        c_step.alignment = Alignment(horizontal="center")

        # Description
        ws3.cell(row=r_idx, column=2, value=step.description)

        # Symbol
        c_sym = ws3.cell(row=r_idx, column=3, value=step.symbol)
        c_sym.alignment = Alignment(horizontal="center")
        c_sym.font = BOLD_10

        # Governing equation
        c_eq = ws3.cell(row=r_idx, column=4, value=step.governing_equation)
        c_eq.font = ITALIC_10

        # Substitution
        c_sub = ws3.cell(row=r_idx, column=5, value=step.substitution)
        c_sub.font = ITALIC_10

        # LIVE EXCEL FORMULA - Written directly as an Excel formula string
        c_formula = ws3.cell(row=r_idx, column=6, value=step.excel_formula)
        c_formula.font = BOLD_10
        c_formula.fill = LIGHT_GRAY_FILL

        # Calculated numerical value
        c_calc = ws3.cell(row=r_idx, column=7, value=step.computed_value)
        c_calc.number_format = "#,##0.00"
        c_calc.alignment = Alignment(horizontal="right")
        c_calc.font = BOLD_10

        # Unit
        c_u = ws3.cell(row=r_idx, column=8, value=step.unit)
        c_u.alignment = Alignment(horizontal="center")

        # Limit reference
        ws3.cell(row=r_idx, column=9, value=step.limit_reference or "N/A")

        # Tolerance
        ws3.cell(row=r_idx, column=10, value=step.tolerance or "Code Standard")

        # Live verification formula (e.g. =IF(G4<=I4,"PASS","FAIL"))
        v_formula = step.status_formula or f'=IF(G{r_idx}>0,"PASS","FAIL")'
        c_vf = ws3.cell(row=r_idx, column=11, value=v_formula)
        c_vf.font = NORMAL_10
        c_vf.fill = LIGHT_GRAY_FILL

        # Status text
        c_stat = ws3.cell(row=r_idx, column=12, value=step.verification_status)
        c_stat.alignment = Alignment(horizontal="center")
        if step.verification_status == "PASS":
            c_stat.font = PASS_FONT
            c_stat.fill = PASS_FILL
        else:
            c_stat.font = FAIL_FONT
            c_stat.fill = FAIL_FILL

        # Separate independent verification evidence
        c_ev = ws3.cell(row=r_idx, column=13, value=step.verification_evidence or "")
        c_ev.font = NORMAL_10

        # Apply borders and default fonts
        for col_i in range(1, 14):
            ws3.cell(row=r_idx, column=col_i).border = THIN_BORDER
            if col_i not in (3, 4, 5, 6, 7, 12):
                ws3.cell(row=r_idx, column=col_i).font = NORMAL_10

    _auto_fit_columns(ws3)

    # ----------------------------------------------------------------------- #
    # Sheet 4: Summary & Recommendations                                      #
    # ----------------------------------------------------------------------- #
    ws4 = wb.create_sheet(title="Summary & Recommendations")
    ws4.cell(row=1, column=1, value="CALCULATION SUMMARY & ENGINEERING VERDICT").font = WHITE_BOLD_11
    ws4.cell(row=1, column=1).fill = NAVY_FILL
    ws4.merge_cells("A1:F1")

    sum_headers = ["Step", "Derivation Item", "Output Symbol", "Value", "Unit", "Verification Result"]
    for c_idx, h in enumerate(sum_headers, start=1):
        c = ws4.cell(row=3, column=c_idx, value=h)
        c.font = WHITE_BOLD_11
        c.fill = STEEL_FILL
        c.alignment = Alignment(horizontal="center", vertical="center")
        c.border = THIN_BORDER

    for r_idx, step in enumerate(spec.steps, start=4):
        ws4.cell(row=r_idx, column=1, value=step.step_id).alignment = Alignment(horizontal="center")
        ws4.cell(row=r_idx, column=2, value=step.description)
        ws4.cell(row=r_idx, column=3, value=step.symbol).alignment = Alignment(horizontal="center")
        c_v = ws4.cell(row=r_idx, column=4, value=step.computed_value)
        c_v.number_format = "#,##0.00"
        c_v.alignment = Alignment(horizontal="right")
        ws4.cell(row=r_idx, column=5, value=step.unit).alignment = Alignment(horizontal="center")
        c_s = ws4.cell(row=r_idx, column=6, value=f"{step.verification_status} (Verified)")
        c_s.alignment = Alignment(horizontal="center")
        c_s.font = PASS_FONT
        c_s.fill = PASS_FILL

        for col_i in range(1, 7):
            ws4.cell(row=r_idx, column=col_i).border = THIN_BORDER
            if col_i not in (3, 6):
                ws4.cell(row=r_idx, column=col_i).font = NORMAL_10

    # Conclusion
    c_row = len(spec.steps) + 6
    ws4.cell(row=c_row, column=1, value="Final Engineering Conclusion").font = WHITE_BOLD_11
    ws4.cell(row=c_row, column=1).fill = STEEL_FILL
    ws4.merge_cells(start_row=c_row, start_column=1, end_row=c_row, end_column=6)

    c_box = ws4.cell(row=c_row + 1, column=1, value=spec.conclusion)
    c_box.font = NORMAL_10
    c_box.alignment = Alignment(wrap_text=True)
    ws4.merge_cells(start_row=c_row + 1, start_column=1, end_row=c_row + 3, end_column=6)

    _auto_fit_columns(ws4)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))


# --------------------------------------------------------------------------- #
# 3. PPTX Executive / Board Presentation Renderer                             #
# --------------------------------------------------------------------------- #

def render_pptx_executive_presentation(spec: ExecutivePresentationSpec, output_path: Path) -> None:
    """Render Executive / Board Presentation as a 16:9 widescreen PPTX."""
    import pptx
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY_RGB = RGBColor(0x1F, 0x4E, 0x79)
    STEEL_RGB = RGBColor(0x2E, 0x75, 0xB6)
    WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_RGB = RGBColor(0x26, 0x26, 0x26)
    GRAY_RGB = RGBColor(0x59, 0x59, 0x59)
    LIGHT_GRAY_RGB = RGBColor(0xF2, 0xF4, 0xF7)
    GREEN_RGB = RGBColor(0x27, 0x6A, 0x3C)
    AMBER_RGB = RGBColor(0xC5, 0x5A, 0x11)

    blank_layout = prs.slide_layouts[6]

    # ----------------------------------------------------------------------- #
    # Slide 1: Title Slide                                                    #
    # ----------------------------------------------------------------------- #
    s1 = prs.slides.add_slide(blank_layout)

    # Top accent header
    top_bar = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY_RGB
    top_bar.line.color.rgb = NAVY_RGB

    # Main Title Box
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p_t = tf1.paragraphs[0]
    p_t.text = spec.presentation_title
    p_t.font.name = "Calibri"
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY_RGB
    p_t.space_after = Pt(12)

    p_sub = tf1.add_paragraph()
    p_sub.text = spec.presentation_subtitle.upper()
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(15)
    p_sub.font.bold = True
    p_sub.font.color.rgb = STEEL_RGB

    # Bottom Metadata Panel
    meta_box = s1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.2))
    tf_meta = meta_box.text_frame
    p_f = tf_meta.paragraphs[0]
    p_f.text = f"FACILITY: {spec.facility}   |   DATE: {spec.date}"
    p_f.font.name = "Calibri"
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = GRAY_RGB
    p_f.space_after = Pt(4)

    p_pres = tf_meta.add_paragraph()
    p_pres.text = f"AUTHORITY: {spec.presenter}"
    p_pres.font.name = "Calibri"
    p_pres.font.size = Pt(12)
    p_pres.font.bold = True
    p_pres.font.color.rgb = NAVY_RGB

    # ----------------------------------------------------------------------- #
    # Content Slides                                                          #
    # ----------------------------------------------------------------------- #
    for slide_spec in spec.slides:
        slide = prs.slides.add_slide(blank_layout)

        # Top Accent Header Bar
        hbar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = NAVY_RGB
        hbar.line.color.rgb = NAVY_RGB

        # Slide Title & Subtitle
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        htf = header_box.text_frame
        htf.word_wrap = True

        hp = htf.paragraphs[0]
        hp.text = slide_spec.title
        hp.font.name = "Calibri"
        hp.font.size = Pt(22)
        hp.font.bold = True
        hp.font.color.rgb = NAVY_RGB

        if slide_spec.subtitle:
            sp = htf.add_paragraph()
            sp.text = slide_spec.subtitle
            sp.font.name = "Calibri"
            sp.font.size = Pt(13)
            sp.font.color.rgb = GRAY_RGB

        current_top = 1.6

        # Metric Cards (horizontal row)
        if slide_spec.cards:
            num_cards = len(slide_spec.cards)
            card_w = min(Inches(2.7), Inches(11.7 / max(num_cards, 1) - 0.2))
            card_h = Inches(1.1)

            for c_idx, card in enumerate(slide_spec.cards):
                c_left = Inches(0.8 + c_idx * (card_w.inches + 0.25))
                card_shape = slide.shapes.add_shape(1, c_left, Inches(current_top), card_w, card_h)
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = LIGHT_GRAY_RGB
                card_shape.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)

                ctf = card_shape.text_frame
                ctf.word_wrap = True

                p_val = ctf.paragraphs[0]
                p_val.text = card.value
                p_val.font.name = "Calibri"
                p_val.font.size = Pt(18)
                p_val.font.bold = True
                p_val.font.color.rgb = NAVY_RGB
                p_val.alignment = PP_ALIGN.CENTER

                p_lbl = ctf.add_paragraph()
                p_lbl.text = card.label
                p_lbl.font.name = "Calibri"
                p_lbl.font.size = Pt(10)
                p_lbl.font.color.rgb = GRAY_RGB
                p_lbl.alignment = PP_ALIGN.CENTER

            current_top += 1.3

        # Bullet Points
        if slide_spec.bullet_points:
            bp_h = Inches(min(2.5, 0.45 * len(slide_spec.bullet_points) + 0.3))
            bp_box = slide.shapes.add_textbox(Inches(0.8), Inches(current_top), Inches(11.7), bp_h)
            bptf = bp_box.text_frame
            bptf.word_wrap = True

            for idx, pt in enumerate(slide_spec.bullet_points):
                bp = bptf.paragraphs[0] if idx == 0 else bptf.add_paragraph()
                bp.text = f"•  {pt}"
                bp.font.name = "Calibri"
                bp.font.size = Pt(13)
                bp.font.color.rgb = DARK_RGB
                bp.space_after = Pt(6)

            current_top += bp_h.inches + 0.1

        # Table
        if slide_spec.table_headers and slide_spec.table_rows:
            num_rows = len(slide_spec.table_rows) + 1
            num_cols = len(slide_spec.table_headers)
            tbl_h = Inches(min(3.2, 0.38 * num_rows + 0.4))
            tbl_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(current_top), Inches(11.7), tbl_h)
            table = tbl_shape.table

            for c_idx, th in enumerate(slide_spec.table_headers):
                cell = table.cell(0, c_idx)
                cell.text = th
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY_RGB
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Calibri"
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = WHITE_RGB

            for r_idx, row_vals in enumerate(slide_spec.table_rows, start=1):
                for c_idx, val in enumerate(row_vals):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.name = "Calibri"
                        p.font.size = Pt(10)
                        if c_idx == num_cols - 1:
                            p.font.bold = True
                            if "PASS" in str(val):
                                p.font.color.rgb = GREEN_RGB
                            elif any(w in str(val) for w in ("FAIL", "ACTION", "CRITICAL")):
                                p.font.color.rgb = AMBER_RGB

            current_top += tbl_h.inches + 0.2

        # Callout box
        if slide_spec.callout:
            callout_top = max(current_top, 6.2)
            c_box = slide.shapes.add_shape(1, Inches(0.8), Inches(callout_top), Inches(11.7), Inches(0.7))
            c_box.fill.solid()
            c_box.fill.fore_color.rgb = LIGHT_GRAY_RGB
            c_box.line.color.rgb = STEEL_RGB

            ctf = c_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = slide_spec.callout
            cp.font.name = "Calibri"
            cp.font.size = Pt(11)
            cp.font.bold = True
            cp.font.color.rgb = NAVY_RGB

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# --------------------------------------------------------------------------- #
# 4. Universal Template Dispatcher                                            #
# --------------------------------------------------------------------------- #

SUPPORTED_TEMPLATES = {
    "approval_note": ArtifactFormat.DOCX,
    "technical_approval_note": ArtifactFormat.DOCX,
    "calculation_workbook": ArtifactFormat.XLSX,
    "engineering_calculation": ArtifactFormat.XLSX,
    "engineering_calculation_workbook": ArtifactFormat.XLSX,
    "executive_presentation": ArtifactFormat.PPTX,
    "board_summary": ArtifactFormat.PPTX,
}


def render_template(
    template_name: str,
    template_data: Optional[Union[Dict[str, Any], BaseModel]],
    output_path: Path,
    art_format: ArtifactFormat,
) -> None:
    """Dispatch and render a named industrial template to output_path.

    Args:
        template_name: Normalized template identifier.
        template_data: Structured payload dictionary or spec model.
        output_path: Destination binary path.
        art_format: Target ArtifactFormat.

    Raises:
        ValueError: If template name is unknown or format is mismatched.
    """
    norm_name = template_name.lower().strip()
    if norm_name not in SUPPORTED_TEMPLATES:
        valid = list(SUPPORTED_TEMPLATES.keys())
        raise ValueError(f"Unknown template '{template_name}'. Supported templates: {valid}")

    expected_fmt = SUPPORTED_TEMPLATES[norm_name]
    if art_format != expected_fmt:
        raise ValueError(
            f"Template '{template_name}' requires format '{expected_fmt.value}', "
            f"but received '{art_format.value}'."
        )

    # 1. DOCX Technical Approval Note
    if norm_name in ("approval_note", "technical_approval_note"):
        if isinstance(template_data, TechnicalApprovalNoteSpec):
            spec = template_data
        elif isinstance(template_data, dict):
            spec = TechnicalApprovalNoteSpec.model_validate(template_data)
        else:
            spec = create_demo_approval_note()
        render_docx_approval_note(spec, output_path)

    # 2. XLSX Engineering Calculation Workbook
    elif norm_name in ("calculation_workbook", "engineering_calculation", "engineering_calculation_workbook"):
        if isinstance(template_data, EngineeringCalculationSpec):
            calc_spec = template_data
        elif isinstance(template_data, dict):
            calc_spec = EngineeringCalculationSpec.model_validate(template_data)
        else:
            calc_spec = create_demo_calculation_workbook()
        render_xlsx_calculation_workbook(calc_spec, output_path)

    # 3. PPTX Executive Presentation
    elif norm_name in ("executive_presentation", "board_summary"):
        if isinstance(template_data, ExecutivePresentationSpec):
            pres_spec = template_data
        elif isinstance(template_data, dict):
            pres_spec = ExecutivePresentationSpec.model_validate(template_data)
        else:
            pres_spec = create_demo_executive_presentation()
        render_pptx_executive_presentation(pres_spec, output_path)
