"""Unit tests for Industrial Deliverables, PPTX Generation, and Template Engine."""

import hashlib
from pathlib import Path
import docx
import openpyxl
from pptx import Presentation
import pytest

from orchestration.capabilities.base import CapabilityContext
from orchestration.capabilities.builtin.artifact import (
    ArtifactFormat,
    ArtifactGenerationCapability,
    ArtifactGenerationRequest,
    PptxGenerator,
)
from orchestration.capabilities.builtin.artifact.templates import (
    EngineeringCalculationSpec,
    ExecutivePresentationSpec,
    TechnicalApprovalNoteSpec,
    create_demo_approval_note,
    create_demo_calculation_workbook,
    create_demo_executive_presentation,
    render_docx_approval_note,
    render_pptx_executive_presentation,
    render_template,
    render_xlsx_calculation_workbook,
)
from orchestration.domain.references import ArtifactReference


@pytest.fixture
def temp_output_dir(tmp_path: Path) -> Path:
    out_dir = tmp_path / "artifacts"
    out_dir.mkdir(parents=True)
    return out_dir


# --------------------------------------------------------------------------- #
# 1. Generic PPTX Generator Tests                                             #
# --------------------------------------------------------------------------- #

def test_pptx_generator_generic(temp_output_dir: Path):
    """Verify generic PptxGenerator creates valid 16:9 presentation from title, content, and data."""
    target = temp_output_dir / "generic_presentation.pptx"
    content = """
## Executive Summary
This is an overview of quarterly plant performance.
- Throughput exceeded target by 4.2%
- Unplanned downtime was zero hours

## Operational Highlights
Operations ran continuously without safety trips.
"""
    table_data = [
        ["Unit", "Availability", "Utilization"],
        ["Crude Unit", "99.8%", "94.2%"],
        ["Hydrocracker", "98.5%", "91.0%"],
    ]

    req = ArtifactGenerationRequest(
        format=ArtifactFormat.PPTX,
        filename="generic_presentation.pptx",
        title="Quarterly Review Presentation",
        content=content,
        data=table_data,
    )
    PptxGenerator.generate(req, target)

    assert target.exists()
    assert target.stat().st_size > 0

    # Inspect with python-pptx
    prs = Presentation(str(target))
    # 16:9 widescreen dimensions: 13.333 x 7.5 inches
    assert round(prs.slide_width.inches, 2) == 13.33
    assert round(prs.slide_height.inches, 2) == 7.50

    # Slides: Title slide + 2 content slides + 1 table slide = 4 slides
    assert len(prs.slides) == 4

    # Slide 1: Title
    s1 = prs.slides[0]
    assert s1.shapes.title is not None
    assert s1.shapes.title.text == "Quarterly Review Presentation"

    # Slide 4: Table slide
    s4 = prs.slides[3]
    table_shapes = [s for s in s4.shapes if s.has_table]
    assert len(table_shapes) == 1
    tbl = table_shapes[0].table
    assert tbl.cell(0, 0).text == "Unit"
    assert tbl.cell(1, 0).text == "Crude Unit"
    assert tbl.cell(1, 1).text == "99.8%"


def test_artifact_capability_pptx_execution(temp_output_dir: Path):
    """Verify ArtifactGenerationCapability produces PPTX with correct MIME type and SHA-256."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-pptx-1")

    res = cap.execute(
        parameters={"artifact_type": "pptx", "filename": "plant_status.pptx", "title": "Plant Status"},
        inputs={"content": "## Section 1\n- Parameter nominal"},
        context=ctx,
    )

    assert res.output is not None
    assert res.output["name"] == "plant_status.pptx"
    assert res.output["mime_type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert len(res.output["sha256"]) == 64
    assert res.output["size_bytes"] > 0

    # Cryptographic integrity check: calculate digest directly from disk
    written_bytes = (temp_output_dir / "plant_status.pptx").read_bytes()
    expected_sha = hashlib.sha256(written_bytes).hexdigest()
    assert res.output["sha256"] == expected_sha

    assert len(res.artifacts) == 1
    art = res.artifacts[0]
    assert isinstance(art, ArtifactReference)
    assert art.name == "plant_status.pptx"
    assert art.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert art.metadata["sha256"] == expected_sha


# --------------------------------------------------------------------------- #
# 2. Technical Approval Note (DOCX) Template Tests                            #
# --------------------------------------------------------------------------- #

def test_docx_technical_approval_note_rendering(temp_output_dir: Path):
    """Verify Technical Approval Note DOCX template renders metadata, tag registry, and sign-offs."""
    target = temp_output_dir / "approval_note.docx"
    spec = create_demo_approval_note()
    render_docx_approval_note(spec, target)

    assert target.exists()

    doc = docx.Document(str(target))
    paragraphs_text = [p.text for p in doc.paragraphs if p.text]

    # Document title and sections present
    assert any("Technical Integrity & Authorization Note" in t for t in paragraphs_text)
    assert any("Executive Summary" in t for t in paragraphs_text)
    assert any("Design Basis & Operating Parameters" in t for t in paragraphs_text)
    assert any("Equipment Inspection & Tag Evaluation Registry" in t for t in paragraphs_text)
    assert any("Engineering Sign-Off & Authority Approval" in t for t in paragraphs_text)

    # Tables inspection:
    # 1. Metadata table, 2. Operating params, 3. Tag registry, 4. Risk assessment, 5. Sign-offs
    assert len(doc.tables) >= 4

    # Tag registry table verification
    tag_table = doc.tables[2]
    header_cells = [c.text for c in tag_table.rows[0].cells]
    assert "Tag ID" in header_cells
    assert "Compliance Status" in header_cells or "Status" in header_cells
    row1_cells = [c.text for c in tag_table.rows[1].cells]
    assert "FV-201A" in row1_cells

    # Sign-off table verification
    sign_off_table = doc.tables[-1]
    so_headers = [c.text for c in sign_off_table.rows[0].cells]
    assert "Role" in so_headers
    assert "Approver Name" in so_headers
    assert "Decision" in so_headers
    so_row1 = [c.text for c in sign_off_table.rows[1].cells]
    assert "A. R. Mitchell, P.E." in so_row1


# --------------------------------------------------------------------------- #
# 3. Engineering Calculation Workbook (XLSX) Template Tests                  #
# --------------------------------------------------------------------------- #

def test_xlsx_engineering_calculation_workbook_rendering(temp_output_dir: Path):
    """Verify Engineering Calculation Workbook XLSX template preserves live formulas and verification evidence."""
    target = temp_output_dir / "calc_workbook.xlsx"
    spec = create_demo_calculation_workbook()
    render_xlsx_calculation_workbook(spec, target)

    assert target.exists()

    # Load with data_only=False to inspect raw Excel formula strings
    wb = openpyxl.load_workbook(str(target), data_only=False)

    # 4 Sheets must exist
    expected_sheets = ["Cover & Scope", "Design Parameters", "Calculations", "Summary & Recommendations"]
    for s_name in expected_sheets:
        assert s_name in wb.sheetnames, f"Expected sheet '{s_name}' missing from workbook"

    # Sheet 1: Verification evidence preserved separately
    ws1 = wb["Cover & Scope"]
    ws1_text = " ".join(str(cell.value or "") for row in ws1.iter_rows() for cell in row)
    assert "Independent Verification Evidence" in ws1_text
    assert "ASME B31.3 Section 304.1.2" in ws1_text

    # Sheet 2: Design Parameters
    ws2 = wb["Design Parameters"]
    assert ws2.cell(row=3, column=2).value == "Parameter Description"
    assert ws2.cell(row=4, column=2).value == "Internal Design Pressure"
    assert ws2.cell(row=4, column=4).value == 4.5

    # Sheet 3: Calculations with live formulas
    ws3 = wb["Calculations"]
    headers = [cell.value for cell in ws3[3]]
    assert "Live Excel Formula" in headers
    assert "Verification Check (Formula)" in headers or "Verification Formula" in headers
    assert "Independent Verification Evidence" in headers

    # Step 1.0: Minimum Pressure Design Wall Thickness
    # Row 4 corresponds to Step 1.0
    formula_val = str(ws3.cell(row=4, column=6).value)
    assert formula_val.startswith("="), f"Expected live formula starting with '=', got: {formula_val}"
    assert "D7" in formula_val  # references pressure

    # Verification formula must be an IF formula starting with =
    v_formula = str(ws3.cell(row=4, column=11).value)
    assert v_formula.startswith("="), f"Expected verification formula starting with '=', got: {v_formula}"
    assert "PASS" in v_formula

    # Step 1.1: Required minimum thickness formula
    formula_1_1 = str(ws3.cell(row=5, column=6).value)
    assert formula_1_1.startswith("=")
    assert "D12" in formula_1_1 or "G8" in formula_1_1

    # Verification status column
    assert ws3.cell(row=4, column=12).value == "PASS"

    # Verification evidence column contains independent audit notes
    evidence_text = str(ws3.cell(row=4, column=13).value)
    assert "Math derivation confirmed" in evidence_text


# --------------------------------------------------------------------------- #
# 4. Executive Presentation (PPTX) Template Tests                             #
# --------------------------------------------------------------------------- #

def test_pptx_executive_presentation_rendering(temp_output_dir: Path):
    """Verify Executive Presentation PPTX template renders 16:9 slides with metric cards and tables."""
    target = temp_output_dir / "board_presentation.pptx"
    spec = create_demo_executive_presentation()
    render_pptx_executive_presentation(spec, target)

    assert target.exists()

    prs = Presentation(str(target))
    assert round(prs.slide_width.inches, 2) == 13.33
    assert round(prs.slide_height.inches, 2) == 7.50

    # Title slide + 5 content slides = 6 slides total
    assert len(prs.slides) == 6

    # Slide 1: Title slide
    s1_text = " ".join(shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame)
    assert "Process Safety & Integrity Review" in s1_text
    assert "Industrial Processing Complex" in s1_text

    # Slide 2: Has metric cards
    s2 = prs.slides[1]
    s2_text = " ".join(shape.text_frame.text for shape in s2.shapes if shape.has_text_frame)
    assert "Executive Summary & Context" in s2_text
    assert "4.5 MPa" in s2_text
    assert "+65.6%" in s2_text

    # Slide 3: Has inspection findings table
    s3 = prs.slides[2]
    table_shapes = [s for s in s3.shapes if s.has_table]
    assert len(table_shapes) >= 1
    t = table_shapes[0].table
    assert t.cell(0, 0).text == "Tag ID"
    assert t.cell(1, 0).text == "FV-201A"

    # Slide 6 (last slide): Decision requested from board
    s_last = prs.slides[5]
    last_text = " ".join(shape.text_frame.text for shape in s_last.shapes if shape.has_text_frame)
    assert "Decision Requested from the Board" in last_text
    assert "$18,500" in last_text


# --------------------------------------------------------------------------- #
# 5. Template Execution via Capability Dispatch                               #
# --------------------------------------------------------------------------- #

def test_capability_template_dispatch_docx(temp_output_dir: Path):
    """Verify capability executes Technical Approval Note template and produces DOCX."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-docx")

    res = cap.execute(
        parameters={"template": "approval_note", "filename": "my_note.docx"},
        inputs={},
        context=ctx,
    )

    assert res.output["name"] == "my_note.docx"
    assert res.output["mime_type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert (temp_output_dir / "my_note.docx").exists()
    doc = docx.Document(str(temp_output_dir / "my_note.docx"))
    assert len(doc.tables) >= 4


def test_capability_template_dispatch_xlsx(temp_output_dir: Path):
    """Verify capability executes Engineering Calculation Workbook template and produces XLSX."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-xlsx")

    res = cap.execute(
        parameters={"template": "calculation_workbook", "filename": "pipe_calc.xlsx"},
        inputs={},
        context=ctx,
    )

    assert res.output["name"] == "pipe_calc.xlsx"
    assert res.output["mime_type"] == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    assert (temp_output_dir / "pipe_calc.xlsx").exists()

    wb = openpyxl.load_workbook(str(temp_output_dir / "pipe_calc.xlsx"), data_only=False)
    assert "Calculations" in wb.sheetnames
    # Check formula in cell F4
    assert str(wb["Calculations"].cell(row=4, column=6).value).startswith("=")


def test_capability_template_dispatch_pptx(temp_output_dir: Path):
    """Verify capability executes Executive Presentation template and produces PPTX."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-pptx")

    res = cap.execute(
        parameters={"template": "executive_presentation", "filename": "board_deck.pptx"},
        inputs={},
        context=ctx,
    )

    assert res.output["name"] == "board_deck.pptx"
    assert res.output["mime_type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert (temp_output_dir / "board_deck.pptx").exists()
    prs = Presentation(str(temp_output_dir / "board_deck.pptx"))
    assert len(prs.slides) == 6


def test_capability_unknown_template_error(temp_output_dir: Path):
    """Verify error raised when invalid template requested."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-err")

    with pytest.raises(ValueError, match="Unknown template 'non_existent'"):
        cap.execute(
            parameters={"template": "non_existent"},
            inputs={},
            context=ctx,
        )


def test_capability_template_format_mismatch_error(temp_output_dir: Path):
    """Verify error raised when requested artifact_type doesn't match template format."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-mismatch")

    with pytest.raises(ValueError, match="Template 'approval_note' requires format 'docx', but received 'pdf'"):
        cap.execute(
            parameters={"artifact_type": "pdf", "template": "approval_note"},
            inputs={},
            context=ctx,
        )


def test_capability_custom_template_data(temp_output_dir: Path):
    """Verify custom template data dictionary overrides default values."""
    cap = ArtifactGenerationCapability(output_dir=temp_output_dir)
    ctx = CapabilityContext(execution_id="exec-tmpl-custom")

    custom_data = {
        "document_id": "CUSTOM-MOC-999",
        "revision": "Rev 2.5",
        "date": "2026-10-15",
        "facility": "Custom Processing Plant Beta",
        "unit_area": "Area 500",
        "title": "Custom High-Temperature Exchanger Authorization",
        "author": "Thermal Systems Engineering",
        "executive_summary": "Custom engineering review of high temperature operating conditions.",
        "design_basis": "ASME Sec VIII Div 1",
        "operating_parameters": {"Design Temperature": "550 °C"},
        "inspection_findings": [
            {
                "tag_id": "E-501",
                "description": "Feed Preheater",
                "pid_reference": "PID-500-01",
                "service": "Preheat",
                "design_spec": "50 barg @ 550 C",
                "measured_condition": "Nominal",
                "compliance_status": "PASS",
            }
        ],
        "recommendations": ["Proceed with caution."],
        "sign_offs": [
            {
                "role": "Lead Engineer",
                "name": "Jane Doe",
                "title": "Principal Engineer",
                "status": "APPROVED",
                "date": "2026-10-15",
            }
        ],
    }

    res = cap.execute(
        parameters={"template": "approval_note", "filename": "custom_note.docx"},
        inputs={"template_data": custom_data},
        context=ctx,
    )

    assert res.output["name"] == "custom_note.docx"
    doc = docx.Document(str(temp_output_dir / "custom_note.docx"))
    text_content = " ".join(p.text for p in doc.paragraphs)
    assert "CUSTOM-MOC-999" in " ".join(cell.text for t in doc.tables for row in t.rows for cell in row.cells)
    assert "Custom High-Temperature Exchanger Authorization" in text_content

